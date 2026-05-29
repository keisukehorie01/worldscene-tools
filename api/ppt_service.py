import base64
import json
import os
import re
import shlex
import shutil
import subprocess
import threading
import time
import uuid
import zipfile
from html import escape
from pathlib import Path
from typing import Any, Dict, Optional

import requests
from flask import jsonify, request, send_file

from billing_sqlite import consume_credit, normalize_email, refund_credit
from drop2ppt_auth import current_auth_email


BASE_DIR = Path(__file__).resolve().parent
RUNTIME_DIR = Path(os.getenv("PPT_RUNTIME_DIR", BASE_DIR / "runtime" / "ppt_jobs"))
UPLOAD_DIR = RUNTIME_DIR / "uploads"
OUTPUT_DIR = RUNTIME_DIR / "outputs"
JOB_DIR = RUNTIME_DIR / "jobs"
CROP_DIR = RUNTIME_DIR / "crops"

MAX_UPLOAD_BYTES = int(os.getenv("PPT_MAX_UPLOAD_BYTES", str(12 * 1024 * 1024)))
MALWARE_SCAN_ENABLED = os.getenv("PPT_MALWARE_SCAN_ENABLED", "true").strip().lower() not in {"0", "false", "no", "off"}
MALWARE_SCAN_REQUIRED = os.getenv("PPT_REQUIRE_MALWARE_SCAN", "true").strip().lower() not in {"0", "false", "no", "off"}
MALWARE_SCAN_COMMAND = os.getenv("PPT_MALWARE_SCAN_COMMAND", "").strip()
MALWARE_SCAN_TIMEOUT = int(os.getenv("PPT_MALWARE_SCAN_TIMEOUT", "60"))
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "").strip()
GEMINI_MODEL = os.getenv("PPT_GEMINI_MODEL", os.getenv("GEMINI_MODEL", "gemini-2.5-flash")).strip()
MAX_IMAGE_REGIONS = int(os.getenv("PPT_MAX_IMAGE_REGIONS", "10"))
STANDARD_MAX_IMAGE_REGIONS = int(os.getenv("PPT_STANDARD_MAX_IMAGE_REGIONS", "4"))
IMAGE_REGION_MIN_AREA = float(os.getenv("PPT_IMAGE_REGION_MIN_AREA", "0.006"))
HIGH_QUALITY_RETRY_MIN_TEXT = int(os.getenv("PPT_HQ_RETRY_MIN_TEXT", "42"))
HIGH_QUALITY_RETRY_MIN_ELEMENTS = int(os.getenv("PPT_HQ_RETRY_MIN_ELEMENTS", "54"))
HIGH_QUALITY_MAX_OUTPUT_TOKENS = int(os.getenv("PPT_HQ_MAX_OUTPUT_TOKENS", "20000"))
STANDARD_MAX_OUTPUT_TOKENS = int(os.getenv("PPT_STANDARD_MAX_OUTPUT_TOKENS", "7000"))
GEMINI_TIMEOUT_SECONDS = int(os.getenv("PPT_GEMINI_TIMEOUT_SECONDS", "120"))
DEFAULT_JA_FONT = os.getenv("PPT_DEFAULT_JA_FONT", "Yu Gothic").strip() or "Yu Gothic"
SERIF_JA_FONT = os.getenv("PPT_SERIF_JA_FONT", "Yu Mincho").strip() or "Yu Mincho"
DEFAULT_LATIN_FONT = os.getenv("PPT_DEFAULT_LATIN_FONT", "Aptos").strip() or "Aptos"

SLIDE_W = 12192000
SLIDE_H = 6858000
BASE_SLIDE_LONG_EDGE = 12192000
JAPANESE_TEXT_RE = re.compile(r"[\u3040-\u30ff\u3400-\u9fff\u3000-\u303f]")

JOBS: Dict[str, Dict[str, Any]] = {}
JOBS_LOCK = threading.Lock()


def normalize_quality(value: Optional[str]) -> str:
    value = (value or "standard").strip().lower()
    if value in {"high_quality", "high-quality", "high"}:
        return "high_quality"
    return "standard"


def credit_type_for_quality(quality: str) -> str:
    return "high_quality" if quality == "high_quality" else "standard"


def register_ppt_routes(app):
    ensure_runtime_dirs()

    @app.route("/api/ppt/jobs", methods=["POST"])
    def create_ppt_job():
      email = current_auth_email()
      if not email:
          return jsonify({"ok": False, "error": "login_required", "message": "Please log in and verify your email."}), 401
      quality = normalize_quality(request.form.get("quality", "standard"))
      credit_type = credit_type_for_quality(quality)

      if "image" not in request.files:
          return jsonify({"ok": False, "error": "image file is required"}), 400

      image = request.files["image"]
      if not image.filename:
          return jsonify({"ok": False, "error": "image filename is required"}), 400

      raw = image.read()
      if not raw:
          return jsonify({"ok": False, "error": "image file is empty"}), 400
      if len(raw) > MAX_UPLOAD_BYTES:
          return jsonify({"ok": False, "error": "image file is too large"}), 413

      mime_type = image.mimetype or "image/png"
      if mime_type not in {"image/png", "image/jpeg", "image/webp"}:
          return jsonify({"ok": False, "error": "PNG, JPG, and WebP are supported"}), 400

      job_id = uuid.uuid4().hex
      suffix = safe_suffix(image.filename, mime_type)
      upload_path = UPLOAD_DIR / f"{job_id}{suffix}"
      output_path = OUTPUT_DIR / f"{job_id}.pptx"
      upload_path.write_bytes(raw)

      job = {
          "id": job_id,
          "status": "queued",
          "progress": 5,
          "message": "Queued",
          "created_at": time.time(),
          "updated_at": time.time(),
          "input_filename": image.filename,
          "input_mime_type": mime_type,
          "input_path": str(upload_path),
          "output_path": str(output_path),
          "email": email,
          "quality": quality,
          "credit_type": credit_type,
          "credits_used": 0,
          "credit_refunded": False,
          "download_url": None,
          "error": None,
      }
      save_job(job)
      threading.Thread(target=process_job, args=(job_id,), daemon=True).start()

      return jsonify({"ok": True, "job": public_job(job)}), 202

    @app.route("/api/ppt/jobs/<job_id>", methods=["GET"])
    def get_ppt_job(job_id: str):
      job = load_job(job_id)
      if not job:
          return jsonify({"ok": False, "error": "job not found"}), 404
      return jsonify({"ok": True, "job": public_job(job)})

    @app.route("/api/ppt/jobs/<job_id>/download", methods=["GET"])
    def download_ppt_job(job_id: str):
      job = load_job(job_id)
      if not job:
          return jsonify({"ok": False, "error": "job not found"}), 404
      if job.get("status") != "completed":
          return jsonify({
              "ok": False,
              "error": "job is not completed",
              "message": "PPTX is still being generated. Please wait a few seconds.",
              "retry_after": 2,
              "job": public_job(job),
          }), 409

      output_path = Path(job["output_path"])
      if not is_output_ready(job):
          try:
              wait_for_output_ready(output_path)
          except RuntimeError:
              return jsonify({
                  "ok": False,
                  "error": "output file is not ready",
                  "message": "PPTX is still being finalized. Please wait a few seconds.",
                  "retry_after": 2,
                  "job": public_job(job),
              }), 409

      return send_file(
          output_path,
          as_attachment=True,
          download_name="drop2ppt-editable.pptx",
          mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
      )


def ensure_runtime_dirs():
    for path in (UPLOAD_DIR, OUTPUT_DIR, JOB_DIR, CROP_DIR):
        path.mkdir(parents=True, exist_ok=True)


def safe_suffix(filename: str, mime_type: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        return suffix
    return {
        "image/png": ".png",
        "image/jpeg": ".jpg",
        "image/webp": ".webp",
    }.get(mime_type, ".img")


def scan_upload_for_malware(path: Path) -> Dict[str, Any]:
    if not MALWARE_SCAN_ENABLED:
        return {"ok": True, "status": "skipped"}

    command = malware_scan_command(path)
    if not command:
        if MALWARE_SCAN_REQUIRED:
            return {
                "ok": False,
                "status": "error",
                "status_code": 503,
                "error": "malware_scanner_unavailable",
                "message": "Malware scanner is not available on the server.",
            }
        return {"ok": True, "status": "skipped"}

    try:
        result = subprocess.run(
            command,
            capture_output=True,
            check=False,
            text=True,
            timeout=MALWARE_SCAN_TIMEOUT,
        )
    except subprocess.TimeoutExpired:
        return {
            "ok": False,
            "status": "error",
            "status_code": 503,
            "error": "malware_scan_timeout",
            "message": "Malware scan timed out.",
        }

    output = "\n".join(part for part in [result.stdout, result.stderr] if part).strip()
    if result.returncode == 0:
        return {"ok": True, "status": "clean", "output": output}
    if result.returncode == 1:
        return {
            "ok": False,
            "status": "infected",
            "status_code": 400,
            "error": "malware_detected",
            "message": "The uploaded file did not pass malware scanning.",
        }
    message = "Malware scanner failed to complete."
    if output:
        message = f"{message} {redact_secret_text(output[:500])}"
    print(
        f"Drop2PPT malware scan failed path={path} returncode={result.returncode} output={redact_secret_text(output[:1000])}",
        flush=True,
    )
    return {
        "ok": False,
        "status": "error",
        "status_code": 503,
        "error": "malware_scan_failed",
        "message": message,
    }


def malware_scan_command(path: Path) -> Optional[list]:
    if MALWARE_SCAN_COMMAND:
        command = shlex.split(MALWARE_SCAN_COMMAND)
        return [part.format(path=str(path)) for part in command]

    clamdscan = shutil.which("clamdscan")
    if clamdscan:
        return [clamdscan, "--no-summary", "--fdpass", str(path)]

    clamscan = shutil.which("clamscan")
    if clamscan:
        return [clamscan, "--no-summary", str(path)]

    return None


def save_job(job: Dict[str, Any]) -> None:
    job["updated_at"] = time.time()
    with JOBS_LOCK:
        JOBS[job["id"]] = job
    path = JOB_DIR / f"{job['id']}.json"
    temp_path = path.with_suffix(".json.tmp")
    temp_path.write_text(json.dumps(job, ensure_ascii=False, indent=2), encoding="utf-8")
    temp_path.replace(path)


def load_job(job_id: str) -> Optional[Dict[str, Any]]:
    path = JOB_DIR / f"{job_id}.json"
    if path.exists():
        try:
            job = json.loads(path.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            job = None
        if job:
            with JOBS_LOCK:
                JOBS[job_id] = job
            return job

    with JOBS_LOCK:
        return JOBS.get(job_id)


def public_job(job: Dict[str, Any]) -> Dict[str, Any]:
    output_ready = is_output_ready(job)
    return {
        "id": job["id"],
        "status": job["status"],
        "progress": job["progress"],
        "message": job["message"],
        "quality": job.get("quality", "standard"),
        "download_url": f"/api/ppt/jobs/{job['id']}/download" if job.get("status") == "completed" and output_ready else None,
        "error": job.get("error"),
    }


def is_output_ready(job: Dict[str, Any]) -> bool:
    output_path = Path(job.get("output_path") or "")
    if job.get("status") != "completed" or not output_path.exists() or output_path.stat().st_size <= 0:
        return False
    try:
        with zipfile.ZipFile(output_path) as archive:
            return archive.testzip() is None
    except (OSError, zipfile.BadZipFile):
        return False


def update_job(job_id: str, **changes) -> Dict[str, Any]:
    job = load_job(job_id)
    if not job:
        raise RuntimeError(f"job not found: {job_id}")
    job.update(changes)
    save_job(job)
    return job


def wait_for_output_ready(output_path: Path, timeout_seconds: float = 20.0) -> None:
    deadline = time.time() + timeout_seconds
    last_error = "output file is not ready"
    while time.time() < deadline:
        if output_path.exists() and output_path.stat().st_size > 0:
            try:
                with zipfile.ZipFile(output_path) as archive:
                    bad_file = archive.testzip()
                if bad_file is None:
                    return
                last_error = f"PPTX archive has a damaged member: {bad_file}"
            except (OSError, zipfile.BadZipFile) as exc:
                last_error = str(exc)
        time.sleep(0.25)
    raise RuntimeError(f"PPTX file was not ready after saving: {last_error}")


def process_job(job_id: str) -> None:
    try:
        job = update_job(job_id, status="processing", progress=10, message="Scanning uploaded image")
        image_path = Path(job["input_path"])

        scan = scan_upload_for_malware(image_path)
        if not scan["ok"]:
            try:
                image_path.unlink(missing_ok=True)
            except OSError:
                pass
            raise RuntimeError(scan.get("message") or scan.get("error") or "Malware scan failed")

        if scan["status"] == "skipped" and MALWARE_SCAN_REQUIRED:
            try:
                image_path.unlink(missing_ok=True)
            except OSError:
                pass
            raise RuntimeError("Malware scanner is required but unavailable.")

        try:
            credit_type = job.get("credit_type") or credit_type_for_quality(job.get("quality", "standard"))
            consume_credit(job["email"], job_id, amount=1, credit_type=credit_type)
            job = update_job(job_id, credits_used=1)
        except ValueError:
            try:
                image_path.unlink(missing_ok=True)
            except OSError:
                pass
            if job.get("quality") == "high_quality":
                raise ValueError("Insufficient High Quality credits. Please buy High Quality before converting.")
            raise ValueError("Insufficient Standard credits. Please buy Starter or Pro before converting.")

        job = update_job(job_id, status="processing", progress=20, message="Analyzing image")
        image_bytes = image_path.read_bytes()
        mime_type = job["input_mime_type"]

        analysis = analyze_image_for_ppt(image_bytes, mime_type, quality=job.get("quality", "standard"))
        analysis["quality"] = job.get("quality", "standard")
        update_job(job_id, progress=70, message="Rebuilding editable slide")

        output_path = Path(job["output_path"])
        write_pptx(output_path, analysis, image_path)
        wait_for_output_ready(output_path)
        update_job(job_id, status="completed", progress=100, message="Ready to download")
    except Exception as exc:
        job = load_job(job_id)
        credits_used = int(job.get("credits_used") or 0) if job else 0
        if job and credits_used > 0 and job.get("email") and not job.get("credit_refunded"):
            refund_credit(
                job["email"],
                job_id,
                amount=int(job.get("credits_used") or 1),
                credit_type=job.get("credit_type") or credit_type_for_quality(job.get("quality", "standard")),
            )
            update_job(job_id, credit_refunded=True)
        error_message = safe_error_message(exc)
        print(f"Drop2PPT job failed job_id={job_id} error={error_message}", flush=True)
        update_job(job_id, status="failed", progress=100, message="Conversion failed", error=error_message)


def safe_error_message(exc: Exception) -> str:
    if isinstance(exc, requests.HTTPError):
        response = exc.response
        status_code = response.status_code if response is not None else "unknown"
        message = ""
        if response is not None:
            try:
                payload = response.json()
                error = payload.get("error") if isinstance(payload, dict) else None
                if isinstance(error, dict):
                    message = clean_text(error.get("message"))
            except ValueError:
                message = clean_text(response.text[:400])
        return redact_secret_text(f"Gemini API error ({status_code}): {message or 'request failed'}")

    if isinstance(exc, json.JSONDecodeError):
        return "AI response was not valid JSON. Please retry the conversion."

    message = redact_secret_text(str(exc))
    return message or f"{type(exc).__name__}: conversion failed"


def redact_secret_text(text: str) -> str:
    text = re.sub(r"([?&]key=)[^&\s]+", r"\1[redacted]", text or "")
    text = re.sub(r"(AIza)[A-Za-z0-9_\-]+", r"\1[redacted]", text)
    text = re.sub(r"(sk|pk|rk)_(live|test)_[A-Za-z0-9_\-]+", r"\1_\2_[redacted]", text)
    return text


def analyze_image_for_ppt(image_bytes: bytes, mime_type: str, quality: str = "standard") -> Dict[str, Any]:
    if not GEMINI_API_KEY:
        return fallback_analysis()

    quality = normalize_quality(quality)
    image_region_instruction = (
        """
Use image_regions only for areas that should remain as bitmap because they are difficult or undesirable
to rebuild as editable objects: photos, realistic illustrations, detailed screenshots that are not meant
to be edited internally, dense decorative textures, QR codes, or highly detailed generated art.
High Quality is not screenshot mode. Never preserve the whole slide as one image.
Do not mark ordinary text cards, buttons, forms, simple boxes, bullets, arrows, logos, simple icons,
or basic diagrams as image_regions; rebuild those as editable PowerPoint text and shapes.
Avoid large image_regions that contain important editable text. Crop only the photographic or
non-editable visual part.
For High Quality output, return many more editable elements than Standard. Use image_regions sparingly:
only for visual-only areas that would clearly look worse if approximated as shapes.
The slide title, headline, subheads, bullets, CTA text, form labels, button captions, plan names,
prices, and table labels must be editable text elements. Missing the main headline is a failure.
Do not collapse lists into one summary. Each visible bullet row, checklist row, icon caption,
comparison-row caption, footer-strip caption, and small CTA caption must be its own editable text
element near the original position.
For landing-page or flyer designs, pay special attention to these text-dense areas:
hero body copy below the headline, benefit icon captions, feature list rows, problem checklist rows,
service/construction captions, before-vs-after comparison bullets, footer feature captions, and final
CTA copy.
For brochure-like Japanese sources, preserve the information architecture: hero headline, body copy,
benefit cards, feature-list table, problem checklist, construction/service gallery captions,
comparison panels, bottom feature bar, and final CTA should all remain editable.
If a photo or phone screenshot contains visible text, preserve the photo/screenshot as an image region
only when that internal text is not expected to be edited. Rebuild all surrounding page text separately.
Coordinates must match the original image location.
"""
        if quality == "high_quality"
        else """
This is a standard hybrid conversion. Rebuild text, panels, arrows, cards, buttons, simple icons,
and layout with editable PowerPoint objects whenever possible. Use image_regions sparingly, up to
4 regions, only for necessary bitmap areas such as photos, logos, QR codes, realistic illustrations,
product images, or detailed screenshots that should not be internally edited.
Do not preserve the whole slide as an image. Avoid image_regions that contain important editable text.
"""
    ).strip()

    prompt = """
You convert visual drafts into editable PowerPoint structure.

Analyze the image and return only valid JSON. Do not wrap it in Markdown.
CRITICAL TEXT FIDELITY RULES:
- Preserve visible text in the original language and wording. Never translate, romanize,
  paraphrase, summarize, or convert Japanese text into English.
- If the image is Japanese, text fields must stay Japanese except for exact visible Latin
  acronyms/words already present in the image, such as AEO, LP, SEO, FAQ, AI, or URL text.
- Do not invent English section labels such as "Hero Section", "Service Strengths",
  "Comparison", or "CTA & Footer" unless those exact English words are visible in the image.
- title/subtitle/sections are only for exact visible source text. If unsure, leave them empty
  and put the visible text in elements instead.
Use this schema:
{
  "title": "short slide title",
  "subtitle": "short slide subtitle",
  "summary": "one sentence answer describing the slide",
  "theme": {
    "background": "hex color",
    "accent": "hex color"
  },
  "elements": [
    {
      "type": "text|rect|roundRect|pill|circle|line",
      "text": "visible text, if any",
      "x": 0.05,
      "y": 0.05,
      "w": 0.20,
      "h": 0.08,
      "fill": "hex color",
      "line": "hex color",
      "font": "hex color",
      "font_face": "Yu Gothic|Yu Mincho|Meiryo|Aptos",
      "font_size": 14,
      "bold": true
    }
  ],
  "image_regions": [
    {
      "purpose": "logo|photo|chart|screenshot|icon_cluster|texture|complex_visual",
      "x": 0.05,
      "y": 0.05,
      "w": 0.20,
      "h": 0.12,
      "layer": "background|foreground",
      "keep_reason": "why this should stay as bitmap"
    }
  ],
  "sections": [
    {
      "title": "section title",
      "body": "1-2 short lines",
      "x": 0.05,
      "y": 0.18,
      "w": 0.28,
      "h": 0.18
    }
  ],
  "steps": [
    {"label": "1", "title": "Upload", "body": "short body"}
  ]
}

Coordinates are normalized 0..1 relative to a 16:9 slide.
For dense infographics, return 30 to 70 elements that preserve the visible layout:
large title, subtitle, numbered cards, side panels, bottom flow blocks, status panels,
thin connectors, progress bars, and simple icon placeholders. Use Japanese text when the image is Japanese.
Prioritize visual placement and editable PowerPoint objects over a generic summary. Do not simplify
the slide into only a few blocks.
Every visible heading, label, button caption, bullet, number, table-like row, and form label must become
an editable text element or a shape with text. Keep x/y/w/h very close to the source image. Adjust each
text box so text does not overlap nearby objects; use smaller font_size and wider boxes when needed.
Keep each text string concise enough to fit its box, but do not move text far from its original position.
For High Quality, normal Japanese flyers, brochures, and landing-page comps should usually contain at
least 42 text-bearing editable elements. Dense infographics or text-heavy LP images should contain 55
or more. If you return fewer than 38 text-bearing elements for a text-heavy Japanese source, the
reconstruction is incomplete. Do not trade editable text for larger bitmap regions.
When a source has stacked rows, comparison boxes, checklist lines, or bottom navigation-like feature
bars, keep each row/caption as a separate editable element instead of merging it into a paragraph.
Set font_face for every text-bearing element. Use common Windows/Office fonts so the PPTX keeps its
appearance on most client PCs: Yu Gothic for Japanese sans/UI text, Yu Mincho for Japanese serif or
elegant poster-like headings, Meiryo for compact Japanese UI labels, and Aptos for Latin text. If unsure,
use Yu Gothic for Japanese and Aptos for Latin. Do not choose decorative or rare fonts.
Do not add explanatory English labels that are not visible in the source image. The output is a
reconstruction of the image, not an analysis report.

{image_region_instruction}
""".strip().replace("{image_region_instruction}", image_region_instruction)

    url = (
        "https://generativelanguage.googleapis.com/v1beta/models/"
        f"{GEMINI_MODEL}:generateContent?key={GEMINI_API_KEY}"
    )

    def request_analysis(request_prompt: str) -> Dict[str, Any]:
        current_prompt = request_prompt
        last_json_error: Optional[json.JSONDecodeError] = None
        for attempt in range(3):
            try:
                return request_analysis_once(current_prompt)
            except json.JSONDecodeError as exc:
                last_json_error = exc
                print(
                    f"Drop2PPT AI JSON parse retry quality={quality} attempt={attempt + 1} error={redact_secret_text(str(exc))}",
                    flush=True,
                )
                if attempt == 0:
                    current_prompt = f"""{request_prompt}

JSON FORMAT RECOVERY:
The previous response was not valid JSON. Return one complete JSON object only.
Do not use Markdown fences, comments, trailing commas, ellipses, or text outside the JSON object.
Keep the structure compact. If the image is very dense, prioritize the main headline, section headings,
large labels, major cards, CTA text, and important list rows so the JSON finishes completely.
"""
                else:
                    current_prompt = f"""{request_prompt}

DENSE REPORT RECOVERY:
The previous response still was not valid JSON, likely because the image is a dense report,
dashboard, score sheet, table-heavy infographic, or chart-heavy document.
Return one complete JSON object only.
Limit the response to at most 45 elements.
For small tables, charts, radar graphs, bar graphs, code blocks, and dense score/detail panels,
use cropped image_regions instead of trying to recreate every tiny label as editable text.
Recreate only the main title, section titles, large score labels, major callouts, and bottom headline
as editable text. The goal is a valid, downloadable, visually faithful PPTX rather than a perfect
line-by-line editable extraction.
Do not include Markdown fences, comments, trailing commas, ellipses, or text outside JSON.
"""
        raise last_json_error or json.JSONDecodeError("AI response did not contain JSON", "", 0)

    def request_analysis_once(request_prompt: str) -> Dict[str, Any]:
        body = {
            "contents": [
                {
                    "parts": [
                        {"text": request_prompt},
                        {
                            "inline_data": {
                                "mime_type": mime_type,
                                "data": base64.b64encode(image_bytes).decode("ascii"),
                            }
                        },
                    ]
                }
            ],
            "generationConfig": {
                "temperature": 0.1 if quality == "high_quality" else 0.2,
                "maxOutputTokens": HIGH_QUALITY_MAX_OUTPUT_TOKENS if quality == "high_quality" else STANDARD_MAX_OUTPUT_TOKENS,
                "responseMimeType": "application/json",
            },
        }
        response = requests.post(url, json=body, timeout=GEMINI_TIMEOUT_SECONDS)
        response.raise_for_status()
        data = response.json()
        candidate = (data.get("candidates") or [{}])[0]
        finish_reason = candidate.get("finishReason")
        parts = ((candidate.get("content") or {}).get("parts")) or []
        text = "".join(part.get("text", "") for part in parts)
        if finish_reason:
            print(f"Drop2PPT Gemini finishReason={redact_secret_text(str(finish_reason))}", flush=True)
        parsed = parse_json_from_model(text)
        return normalize_analysis(parsed)

    analysis = request_analysis(prompt)
    if quality == "high_quality" and high_quality_needs_retry(analysis):
        retry_prompt = f"""{prompt}

STRICT RETRY FOR HIGH QUALITY:
Your previous structure likely preserved too much as bitmap or missed visible text.
Return a fuller editable reconstruction. Include every visible Japanese headline, subhead, bullet,
CTA, caption, form label, price, service name, plan name, and comparison label as an editable
text-bearing element. Use bitmap image_regions only for photo/realistic/screenshot content that
would clearly be worse as PowerPoint shapes. Do not include surrounding editable text inside those
image regions. Prefer smaller cropped image regions plus editable text placed over or beside them.
Keep Japanese text in Japanese. Do not return English summary section names or English descriptions
unless those exact English words are visible in the source image.
For LP/flyer images, explicitly include: hero body copy, benefit icon captions, every feature-list row,
every checklist row, each service/construction caption, each comparison bullet, footer strip captions,
and final CTA captions.
The retry result should have at least {HIGH_QUALITY_RETRY_MIN_TEXT} text-bearing editable elements
unless the source image is truly sparse.
"""
        try:
            retry_analysis = request_analysis(retry_prompt)
            if high_quality_analysis_score(retry_analysis) >= high_quality_analysis_score(analysis):
                analysis = retry_analysis
        except Exception as exc:
            print(f"Drop2PPT high quality retry skipped error={redact_secret_text(str(exc))}", flush=True)
    if quality != "high_quality":
        analysis["image_regions"] = (analysis.get("image_regions") or [])[:STANDARD_MAX_IMAGE_REGIONS]
    return analysis


def parse_json_from_model(text: str) -> Dict[str, Any]:
    cleaned = (text or "").strip().lstrip("\ufeff")
    last_error: Optional[json.JSONDecodeError] = None

    for candidate in json_candidates_from_model_text(cleaned):
        try:
            parsed = json.loads(candidate)
        except json.JSONDecodeError as exc:
            last_error = exc
            continue
        if not isinstance(parsed, dict):
            last_error = json.JSONDecodeError("Top-level JSON value must be an object", candidate, 0)
            continue
        return parsed

    raise last_error or json.JSONDecodeError("No complete JSON object found", cleaned, 0)


def json_candidates_from_model_text(text: str):
    seen = set()

    def add(candidate: str):
        candidate = strip_markdown_json_fence(candidate.strip())
        if candidate and candidate not in seen:
            seen.add(candidate)
            return candidate
        return None

    direct = add(text)
    if direct:
        yield direct

    for match in re.finditer(r"```(?:json)?\s*([\s\S]*?)```", text, flags=re.IGNORECASE):
        candidate = add(match.group(1))
        if candidate:
            yield candidate

    balanced = extract_balanced_json_objects(text)
    balanced.sort(key=len, reverse=True)
    for candidate_text in balanced:
        candidate = add(candidate_text)
        if candidate:
            yield candidate


def strip_markdown_json_fence(text: str) -> str:
    text = re.sub(r"^\s*```(?:json)?\s*", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\s*```\s*$", "", text)
    return text.strip()


def extract_balanced_json_objects(text: str):
    objects = []
    start = None
    depth = 0
    in_string = False
    escaped = False

    for index, char in enumerate(text):
        if in_string:
            if escaped:
                escaped = False
            elif char == "\\":
                escaped = True
            elif char == '"':
                in_string = False
            continue

        if char == '"':
            in_string = True
            continue
        if char == "{":
            if depth == 0:
                start = index
            depth += 1
            continue
        if char == "}" and depth:
            depth -= 1
            if depth == 0 and start is not None:
                objects.append(text[start:index + 1])
                start = None

    return objects


def high_quality_needs_retry(analysis: Dict[str, Any]) -> bool:
    text_count = text_bearing_element_count(analysis)
    element_count = len(analysis.get("elements") or [])
    regions = analysis.get("image_regions") or []
    large_regions = sum(1 for region in regions if image_region_area(region) > 0.18)
    total_region_area = sum(image_region_area(region) for region in regions)
    return (
        text_count < HIGH_QUALITY_RETRY_MIN_TEXT
        or element_count < HIGH_QUALITY_RETRY_MIN_ELEMENTS
        or large_regions >= 3
        or (total_region_area > 0.48 and text_count < HIGH_QUALITY_RETRY_MIN_TEXT + 8)
    )


def high_quality_analysis_score(analysis: Dict[str, Any]) -> float:
    text_count = text_bearing_element_count(analysis)
    element_count = len(analysis.get("elements") or [])
    regions = analysis.get("image_regions") or []
    large_region_penalty = sum(10 for region in regions if image_region_area(region) > 0.18)
    area_penalty = sum(image_region_area(region) for region in regions) * 20
    return text_count * 5 + element_count - large_region_penalty - area_penalty


def text_bearing_element_count(analysis: Dict[str, Any]) -> int:
    count = 0
    for element in analysis.get("elements") or []:
        if isinstance(element, dict) and str(element.get("text") or "").strip():
            count += 1
    return count


def image_region_area(region: Any) -> float:
    if not isinstance(region, dict):
        return 0.0
    try:
        width = float(region.get("w") or 0)
        height = float(region.get("h") or 0)
    except (TypeError, ValueError):
        return 0.0
    return max(0.0, width) * max(0.0, height)


def normalize_analysis(raw: Dict[str, Any]) -> Dict[str, Any]:
    analysis = fallback_analysis()
    analysis["title"] = clean_text(raw.get("title")) or analysis["title"]
    analysis["subtitle"] = clean_text(raw.get("subtitle")) or analysis["subtitle"]
    analysis["summary"] = clean_text(raw.get("summary")) or analysis["summary"]

    theme = raw.get("theme") if isinstance(raw.get("theme"), dict) else {}
    analysis["theme"] = {
        "background": clean_hex(theme.get("background")) or analysis["theme"]["background"],
        "accent": clean_hex(theme.get("accent")) or analysis["theme"]["accent"],
    }

    elements = []
    for item in raw.get("elements") or []:
        if not isinstance(item, dict):
            continue
        element_type = str(item.get("type") or "rect").strip()
        element_type = "roundRect" if element_type.lower() in {"roundrect", "rounded_rect", "rounded"} else element_type.lower()
        if element_type not in {"text", "rect", "roundRect", "pill", "circle", "line"}:
            element_type = "rect"
        default_fill = "10233F" if element_type in {"text", "line"} else "123B59"
        elements.append({
            "type": element_type,
            "text": clean_text(item.get("text")),
            "x": clamp_float(item.get("x"), 0.0, 0.98),
            "y": clamp_float(item.get("y"), 0.0, 0.98),
            "w": clamp_float(item.get("w"), 0.0 if element_type == "line" else 0.01, 1.0),
            "h": clamp_float(item.get("h"), 0.0 if element_type == "line" else 0.01, 1.0),
            "fill": clean_hex(item.get("fill")) or default_fill,
            "line": clean_hex(item.get("line")) or analysis["theme"]["accent"],
            "font": clean_hex(item.get("font")) or "FFFFFF",
            "font_face": clean_font_face(item.get("font_face") or item.get("font_family") or item.get("typeface"), clean_text(item.get("text")), item.get("font_size")),
            "font_size": clamp_int(item.get("font_size"), 7, 44),
            "bold": bool(item.get("bold", element_type in {"text", "pill", "circle"})),
        })
    if elements:
        analysis["elements"] = elements[:120]

    image_regions = []
    raw_regions = raw.get("image_regions") or raw.get("bitmap_regions") or raw.get("preserve_regions") or []
    for item in raw_regions:
        if not isinstance(item, dict):
            continue
        x = clamp_float(item.get("x"), 0.0, 0.98)
        y = clamp_float(item.get("y"), 0.0, 0.98)
        w = clamp_float(item.get("w"), 0.01, 1.0)
        h = clamp_float(item.get("h"), 0.01, 1.0)
        if x + w > 1.0:
            w = max(0.01, 1.0 - x)
        if y + h > 1.0:
            h = max(0.01, 1.0 - y)
        if w * h < IMAGE_REGION_MIN_AREA:
            continue
        layer = clean_text(item.get("layer")).lower()
        if layer not in {"background", "foreground"}:
            layer = "foreground"
        image_regions.append({
            "purpose": clean_text(item.get("purpose"))[:40] or "complex_visual",
            "x": x,
            "y": y,
            "w": w,
            "h": h,
            "layer": layer,
            "keep_reason": clean_text(item.get("keep_reason"))[:160],
        })
    if image_regions:
        analysis["image_regions"] = image_regions[:MAX_IMAGE_REGIONS]

    sections = []
    for item in raw.get("sections") or []:
        if not isinstance(item, dict):
            continue
        sections.append({
            "title": clean_text(item.get("title")) or "Section",
            "body": clean_text(item.get("body")) or "",
            "x": clamp_float(item.get("x"), 0.05, 0.86),
            "y": clamp_float(item.get("y"), 0.18, 0.78),
            "w": clamp_float(item.get("w"), 0.18, 0.40),
            "h": clamp_float(item.get("h"), 0.12, 0.22),
        })
    if sections:
        analysis["sections"] = sections[:10]

    steps = []
    for item in raw.get("steps") or []:
        if not isinstance(item, dict):
            continue
        steps.append({
            "label": clean_text(item.get("label")) or str(len(steps) + 1),
            "title": clean_text(item.get("title")) or "Step",
            "body": clean_text(item.get("body")) or "",
        })
    if steps:
        analysis["steps"] = steps[:5]

    return analysis


def fallback_analysis() -> Dict[str, Any]:
    return {
        "title": "Image to PowerPoint",
        "subtitle": "The source image is preserved because AI analysis was unavailable.",
        "summary": "The uploaded image is kept on the slide so it can still be reviewed and rebuilt manually.",
        "theme": {
            "background": "F6FAFC",
            "accent": "0E8F82",
        },
        "elements": fallback_elements(),
        "image_regions": [{"x": 0.0, "y": 0.0, "w": 1.0, "h": 1.0, "layer": "background"}],
        "sections": [
            {"title": "Source image", "body": "Preserved for visual fidelity", "x": 0.05, "y": 0.78, "w": 0.28, "h": 0.10},
            {"title": "Manual cleanup", "body": "Run conversion again for editable reconstruction", "x": 0.36, "y": 0.78, "w": 0.30, "h": 0.10},
            {"title": "PPTX", "body": "Generated fallback slide", "x": 0.69, "y": 0.78, "w": 0.26, "h": 0.10},
        ],
        "steps": [
            {"label": "1", "title": "Upload", "body": "Drop an image"},
            {"label": "2", "title": "Analyze", "body": "Read the layout"},
            {"label": "3", "title": "Preserve", "body": "Keep the original visual"},
            {"label": "4", "title": "Download", "body": "Get a PPTX file"},
        ],
    }


def fallback_elements():
    return [
        {"type": "rect", "x": 0.03, "y": 0.03, "w": 0.94, "h": 0.08, "fill": "FFFFFF", "line": "0E8F82"},
        {"type": "text", "text": "Image preserved because AI analysis was unavailable", "x": 0.06, "y": 0.052, "w": 0.70, "h": 0.035, "font": "0B2B40", "font_size": 18, "bold": True},
        {"type": "text", "text": "Please retry the conversion if you need more editable objects.", "x": 0.06, "y": 0.88, "w": 0.70, "h": 0.035, "font": "536678", "font_size": 12},
    ]

def build_dense_reconstruction(analysis: Dict[str, Any]):
    cards = extract_step_cards(analysis)
    logs = extract_log_lines(analysis)
    title = clean_text(analysis.get("title")) or "Image to PowerPoint"
    subtitle = clean_text(analysis.get("subtitle")) or "Editable reconstruction"
    summary = clean_text(analysis.get("summary")) or "Generated from the uploaded image."

    elements = []
    add_bg_grid(elements)
    elements.append({"type": "imageRegions", "layer": "background"})
    add_text(elements, title, 0.23, 0.028, 0.54, 0.075, 34, bold=True)
    add_text(elements, subtitle, 0.25, 0.112, 0.54, 0.035, 13, font="EAF7FF", bold=True)

    add_left_log_panel(elements, logs)
    add_center_core(elements)
    add_process_cards(elements, cards)
    add_status_panel(elements)
    add_bottom_flow(elements)
    add_text(elements, summary, 0.10, 0.916, 0.80, 0.045, 14, font="FFFFFF", bold=True)
    return elements


def extract_step_cards(analysis: Dict[str, Any]):
    defaults = default_step_cards()
    cards = {}
    for step in analysis.get("steps") or []:
        if not isinstance(step, dict):
            continue
        label = clean_text(step.get("label")) or str(len(cards) + 1)
        parsed = parse_step_heading(label)
        if parsed:
            label, parsed_title = parsed
        else:
            parsed_title = ""
        title = clean_text(step.get("title")) or parsed_title or defaults.get(label, {}).get("title", "Step")
        body = clean_text(step.get("body")) or defaults.get(label, {}).get("body", "")
        if label.isdigit():
            cards[label] = {"label": label, "title": title, "body": body}

    for element in analysis.get("elements") or []:
        text = clean_text(element.get("text") if isinstance(element, dict) else "")
        if not text:
            continue
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        if not lines:
            continue
        match = re.match(r"^(\d{1,2})\s*[.:\-)]?\s*(.+)$", lines[0])
        if not match:
            continue
        label, heading = match.groups()
        number = int(label)
        if not 1 <= number <= 12:
            continue
        body = "\n".join(lines[1:3])
        if not body:
            body = defaults.get(label, {}).get("body", "")
        cards[label] = {"label": label, "title": heading[:28], "body": body[:90]}

    for label, data in defaults.items():
        cards.setdefault(label, data)
    return [cards[str(index)] for index in range(1, 13)]


def default_step_cards():
    rows = [
        ("1", "Upload", "Receive the uploaded image"),
        ("2", "Scan", "Check the file before processing"),
        ("3", "Analyze", "Read the visible layout and text"),
        ("4", "Structure", "Group important visual sections"),
        ("5", "Extract text", "Recreate readable text where possible"),
        ("6", "Preserve images", "Keep photo-like areas as image regions"),
        ("7", "Rebuild", "Create editable PowerPoint objects"),
        ("8", "Style", "Apply matching colors and typography"),
        ("9", "Review", "Check basic slide consistency"),
        ("10", "Package", "Write the generated PPTX file"),
        ("11", "Finalize", "Prepare the download file"),
        ("12", "Download", "Return the finished PowerPoint"),
    ]
    return {label: {"label": label, "title": title, "body": body} for label, title, body in rows}

def extract_log_lines(analysis: Dict[str, Any]):
    defaults = []
    found = []
    for element in analysis.get("elements") or []:
        if not isinstance(element, dict):
            continue
        for line in clean_text(element.get("text")).splitlines():
            line = line.strip()
            if ">" in line and len(line) >= 5:
                found.append(line[:42])
    merged = []
    for line in found + defaults:
        if line not in merged:
            merged.append(line)
    return merged[:12]


def add_bg_grid(elements):
    elements.append({"type": "rect", "x": 0, "y": 0, "w": 1, "h": 1, "fill": "06162A", "line": "06162A"})
    for index in range(18):
        y = 0.05 + index * 0.048
        add_line(elements, 0.01, y, 0.97, 0, line="082D52")
    for index in range(28):
        x = 0.02 + index * 0.035
        add_line(elements, x, 0.02, 0, 0.88, line="071F3D")
    for index in range(32):
        x = (index * 37 % 96) / 100
        y = (index * 53 % 88) / 100 + 0.02
        color = "21D7FF" if index % 3 else "FFB23E"
        elements.append({"type": "circle", "x": x, "y": y, "w": 0.0035, "h": 0.006, "fill": color, "line": color})
    for index in range(22):
        x = 0.01 + (index % 11) * 0.085
        y = 0.02 + (index // 11) * 0.89
        add_line(elements, x, y, 0.045, 0, line="19B9FF")
        add_line(elements, x + 0.045, y, 0.018, 0.025, line="19B9FF")


def add_left_log_panel(elements, logs):
    add_panel(elements, 0.01, 0.09, 0.16, 0.56, "Processing log", line="0CE879", fill="021A24")
    for index, line in enumerate(logs):
        y = 0.145 + index * 0.0375
        add_text(elements, line, 0.012, y, 0.145, 0.021, 7, font="33F28A", bold=True)
        add_line(elements, 0.012, y + 0.026, 0.135, 0, line="0D7C55")
    for index in range(16):
        height = 0.012 + (index % 5) * 0.006
        elements.append({"type": "rect", "x": 0.018 + index * 0.0084, "y": 0.625 - height, "w": 0.005, "h": height, "fill": "25F075", "line": "25F075"})


def add_center_core(elements):
    for size, color in [(0.30, "0B3D70"), (0.24, "116FB2"), (0.18, "21D7FF"), (0.12, "FFB23E")]:
        x = 0.50 - size / 2
        y = 0.39 - size / 3.2
        elements.append({"type": "circle", "x": x, "y": y, "w": size, "h": size * 0.56, "fill": "06162A", "line": color})
    elements.append({"type": "roundRect", "x": 0.44, "y": 0.315, "w": 0.12, "h": 0.16, "fill": "0A5FA8", "line": "7FE6FF"})
    add_text(elements, "AI", 0.462, 0.354, 0.076, 0.075, 36, bold=True)
    for index in range(18):
        x = 0.37 + (index % 6) * 0.046
        y = 0.255 + (index // 6) * 0.065
        elements.append({"type": "rect", "x": x, "y": y, "w": 0.018, "h": 0.010, "fill": "103B64", "line": "21D7FF"})
    for x, y, w, h in [(0.18, 0.19, 0.18, 0.16), (0.64, 0.18, -0.08, 0.12), (0.19, 0.58, 0.20, -0.12), (0.62, 0.58, -0.12, -0.10)]:
        add_line(elements, x, y, w, h, line="38D5FF")


def add_process_cards(elements, cards):
    positions = {
        1: (0.19, 0.15), 2: (0.32, 0.15), 3: (0.48, 0.15), 4: (0.64, 0.15),
        5: (0.65, 0.29), 6: (0.65, 0.43), 7: (0.61, 0.57), 8: (0.45, 0.57),
        9: (0.31, 0.57), 10: (0.17, 0.57), 11: (0.17, 0.43), 12: (0.17, 0.29),
    }
    for card in cards:
        label = int(card["label"])
        x, y = positions.get(label, (0.2, 0.2))
        w, h = 0.145, 0.105
        add_infographic_card(elements, x, y, w, h, card["label"], card["title"], card["body"])
    for x, y, w, h in [
        (0.292, 0.196, 0.028, 0), (0.455, 0.196, 0.025, 0), (0.615, 0.196, 0.025, 0),
        (0.720, 0.255, 0.030, 0.045), (0.720, 0.395, 0.030, 0.045),
        (0.588, 0.615, -0.032, 0), (0.435, 0.615, -0.030, 0), (0.293, 0.615, -0.030, 0),
        (0.170, 0.535, 0, -0.035), (0.170, 0.395, 0, -0.035),
    ]:
        add_line(elements, x, y, w, h, line="38D5FF")


def add_infographic_card(elements, x, y, w, h, label, title, body):
    elements.append({"type": "roundRect", "x": x, "y": y, "w": w, "h": h, "fill": "08233E", "line": "38D5FF"})
    elements.append({"type": "rect", "x": x + 0.006, "y": y + 0.010, "w": 0.035, "h": h - 0.020, "fill": "0B2D50", "line": "176EA8"})
    elements.append({"type": "circle", "text": label, "x": x - 0.010, "y": y - 0.007, "w": 0.030, "h": 0.040, "fill": "0B6EEA", "line": "7FE6FF", "font_size": 13, "bold": True})
    add_icon_placeholder(elements, x + 0.014, y + 0.032, label)
    add_text(elements, title, x + 0.048, y + 0.018, w - 0.055, 0.030, 13, bold=True)
    for index, line in enumerate(wrap_short(body, 16)[:3]):
        add_text(elements, line, x + 0.050, y + 0.050 + index * 0.020, w - 0.057, 0.018, 7, font="EAF7FF")
    add_line(elements, x + 0.012, y + h - 0.012, w - 0.024, 0, line="1BB8FF")


def add_icon_placeholder(elements, x, y, seed):
    elements.append({"type": "circle", "x": x, "y": y, "w": 0.020, "h": 0.030, "fill": "073A65", "line": "7FE6FF"})
    add_line(elements, x + 0.004, y + 0.015, 0.012, -0.010, line="7FE6FF")
    add_line(elements, x + 0.004, y + 0.015, 0.012, 0.010, line="7FE6FF")
    if int(seed) % 2 == 0:
        elements.append({"type": "rect", "x": x + 0.024, "y": y + 0.006, "w": 0.014, "h": 0.018, "fill": "093A62", "line": "7FE6FF"})
    else:
        elements.append({"type": "circle", "x": x + 0.024, "y": y + 0.006, "w": 0.014, "h": 0.018, "fill": "093A62", "line": "7FE6FF"})


def add_status_panel(elements):
    add_panel(elements, 0.82, 0.03, 0.165, 0.67, "処理ステータス", fill="071A33", line="1AA6D9")
    add_text(elements, "Stage", 0.845, 0.095, 0.050, 0.035, 14, bold=True)
    add_text(elements, "10", 0.895, 0.084, 0.040, 0.050, 24, font="66F0E8", bold=True)
    add_text(elements, "/ 12", 0.933, 0.100, 0.036, 0.030, 13, bold=True)
    add_text(elements, "カンプ解析中", 0.852, 0.145, 0.105, 0.035, 15, font="00E7FF", bold=True)
    for size, color in [(0.110, "10507E"), (0.088, "22D7FF"), (0.066, "66E95E")]:
        elements.append({"type": "circle", "x": 0.855 + (0.110 - size) / 2, "y": 0.205 + (0.110 - size) / 2, "w": size, "h": size, "fill": "071A33", "line": color})
    add_text(elements, "87%", 0.884, 0.245, 0.060, 0.045, 24, bold=True)
    add_text(elements, "進行状況", 0.828, 0.345, 0.060, 0.020, 8, bold=True)
    for index in range(12):
        color = "37EB78" if index < 10 else "0D3557"
        elements.append({"type": "roundRect", "x": 0.829 + index * 0.0118, "y": 0.374, "w": 0.009, "h": 0.012, "fill": color, "line": color})
    elements.append({"type": "rect", "x": 0.828, "y": 0.405, "w": 0.145, "h": 0.105, "fill": "061D35", "line": "1AA6D9"})
    add_text(elements, "現在の処理", 0.835, 0.414, 0.070, 0.020, 9, font="FFFFFF", bold=True)
    for index, text in enumerate(["・カンプ画像を解析中...", "・サービスカードを抽出", "・FAQを認識", "・CTAを構造化しています..."]):
        add_text(elements, text, 0.838, 0.441 + index * 0.017, 0.126, 0.017, 7, font="EAF7FF")
    elements.append({"type": "rect", "x": 0.828, "y": 0.530, "w": 0.145, "h": 0.145, "fill": "061D35", "line": "1AA6D9"})
    add_text(elements, "最近のログ", 0.835, 0.538, 0.080, 0.020, 9, bold=True)
    for index in range(6):
        add_text(elements, f"12:3{index}:1{index}  処理ステップを完了しました", 0.837, 0.565 + index * 0.017, 0.126, 0.016, 6, font="CFEFFF")
        elements.append({"type": "circle", "x": 0.962, "y": 0.568 + index * 0.017, "w": 0.006, "h": 0.008, "fill": "39EF78", "line": "39EF78"})


def add_bottom_flow(elements):
    add_panel(elements, 0.01, 0.70, 0.16, 0.20, "重要キーワード", fill="062646", line="25BFFF")
    keywords = ["個別設計", "テンプレ禁止", "カンプ解析", "HTML分解", "AEO対策", "LLMO対策", "問い合わせ導線", "スマホ最適化"]
    for index, text in enumerate(keywords):
        x = 0.020 + (index % 2) * 0.072
        y = 0.745 + (index // 2) * 0.035
        elements.append({"type": "roundRect", "x": x, "y": y, "w": 0.064, "h": 0.024, "fill": "073A65", "line": "25BFFF"})
        add_text(elements, text, x + 0.006, y + 0.005, 0.052, 0.014, 6, bold=True)

    panels = [
        ("1 デザインカンプ生成", 0.19, 0.70, 0.27),
        ("2 カンプ解析", 0.48, 0.70, 0.16),
        ("3 HTML分解", 0.66, 0.70, 0.15),
        ("4 個別LP完成", 0.83, 0.70, 0.15),
    ]
    for title, x, y, w in panels:
        add_panel(elements, x, y, w, 0.20, title, fill="062646", line="25BFFF")
    for index, color in enumerate(["0E56A0", "47A953", "42AEEB", "E49A20"]):
        x = 0.205 + index * 0.065
        elements.append({"type": "rect", "x": x, "y": 0.748, "w": 0.050, "h": 0.102, "fill": color, "line": "D8F6FF"})
        add_text(elements, f"{chr(65 + index)}案", x + 0.012, 0.756, 0.028, 0.020, 9, bold=True)
        for line_index in range(4):
            elements.append({"type": "rect", "x": x + 0.007, "y": 0.785 + line_index * 0.014, "w": 0.036, "h": 0.006, "fill": "EAF7FF", "line": "EAF7FF"})
    for index, text in enumerate(["ヒーロー領域", "サービスカード", "導入フロー", "FAQセクション", "CTAボタン"]):
        add_text(elements, text, 0.555, 0.747 + index * 0.026, 0.070, 0.016, 7, font="FFFFFF")
        elements.append({"type": "circle", "x": 0.542, "y": 0.750 + index * 0.026, "w": 0.007, "h": 0.009, "fill": "FFD54D", "line": "FFD54D"})
    code_lines = ["<header class=\"hero\">", "<section class=\"cards\">", "<section class=\"flow\">", "<section class=\"faq\">", "<footer class=\"contact\">"]
    for index, text in enumerate(code_lines):
        add_text(elements, text, 0.680, 0.748 + index * 0.026, 0.115, 0.018, 7, font="DDF6FF")
    elements.append({"type": "rect", "x": 0.845, "y": 0.758, "w": 0.075, "h": 0.080, "fill": "EAF7FF", "line": "25BFFF"})
    elements.append({"type": "rect", "x": 0.925, "y": 0.765, "w": 0.035, "h": 0.080, "fill": "EAF7FF", "line": "25BFFF"})
    for index in range(4):
        elements.append({"type": "rect", "x": 0.852, "y": 0.775 + index * 0.014, "w": 0.060, "h": 0.006, "fill": "0E56A0", "line": "0E56A0"})
        elements.append({"type": "rect", "x": 0.931, "y": 0.782 + index * 0.013, "w": 0.022, "h": 0.005, "fill": "0E56A0", "line": "0E56A0"})
    elements.append({"type": "circle", "text": "✓", "x": 0.950, "y": 0.828, "w": 0.035, "h": 0.050, "fill": "42C946", "line": "D5FFD8", "font_size": 20, "bold": True})
    for x in [0.455, 0.640, 0.810]:
        add_line(elements, x, 0.800, 0.020, 0, line="7FE6FF")
        add_line(elements, x + 0.020, 0.800, -0.008, -0.010, line="7FE6FF")
        add_line(elements, x + 0.020, 0.800, -0.008, 0.010, line="7FE6FF")


def add_panel(elements, x, y, w, h, title, fill="071F3D", line="1AA6D9"):
    elements.append({"type": "roundRect", "x": x, "y": y, "w": w, "h": h, "fill": fill, "line": line})
    add_line(elements, x + 0.010, y + 0.036, w - 0.020, 0, line=line)
    add_text(elements, title, x + 0.016, y + 0.013, w - 0.030, 0.022, 10, font="FFFFFF", bold=True)
    for dx, dy in [(0.004, 0.004), (w - 0.018, 0.004), (0.004, h - 0.018), (w - 0.018, h - 0.018)]:
        elements.append({"type": "rect", "x": x + dx, "y": y + dy, "w": 0.014, "h": 0.004, "fill": line, "line": line})


def add_text(elements, text, x, y, w, h, size, font="FFFFFF", bold=False):
    elements.append({"type": "text", "text": text, "x": x, "y": y, "w": w, "h": h, "font": font, "font_size": size, "bold": bold})


def add_line(elements, x, y, w, h, line="1AA6D9"):
    elements.append({"type": "line", "x": x, "y": y, "w": w, "h": h, "line": line})


def wrap_short(text: str, size: int):
    text = clean_text(text)
    if not text:
        return []
    if "\n" in text:
        return [line for line in text.splitlines() if line]
    return [text[index:index + size] for index in range(0, len(text), size)]


def clean_text(value: Any) -> str:
    if value is None:
        return ""
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in str(value).splitlines()]
    text = "\n".join(line for line in lines if line).strip()
    return text[:180]


def clean_hex(value: Any) -> str:
    if value is None:
        return ""
    text = str(value).strip().lstrip("#").upper()
    return text if re.fullmatch(r"[0-9A-F]{6}", text) else ""


def clean_font_face(value: Any, text: str = "", font_size: Any = None) -> str:
    raw = clean_text(value).replace("\n", " ").strip()
    lowered = raw.lower()
    if raw:
        if "mincho" in lowered or "serif" in lowered or "明朝" in raw:
            return SERIF_JA_FONT
        if "meiryo" in lowered or "メイリオ" in raw:
            return "Meiryo"
        if "gothic" in lowered or "sans" in lowered or "ゴシック" in raw:
            return DEFAULT_JA_FONT
        if "aptos" in lowered:
            return DEFAULT_LATIN_FONT
        if raw in {DEFAULT_JA_FONT, SERIF_JA_FONT, DEFAULT_LATIN_FONT, "Meiryo", "Yu Gothic", "Yu Mincho", "Aptos"}:
            return raw[:60]
    return default_font_face(text, font_size)


def default_font_face(text: str = "", font_size: Any = None) -> str:
    size = clamp_int(font_size, 0, 200)
    if JAPANESE_TEXT_RE.search(text or ""):
        # Large Japanese poster headings often look closer to Mincho; UI/body text is safer in Gothic.
        return SERIF_JA_FONT if size >= 30 else DEFAULT_JA_FONT
    return DEFAULT_LATIN_FONT


def apply_run_font(run, font_face: str) -> None:
    font_face = clean_font_face(font_face)
    if not font_face:
        return
    run.font.name = font_face
    try:
        from pptx.oxml import OxmlElement
        from pptx.oxml.ns import qn

        rpr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            node = rpr.find(qn(tag))
            if node is None:
                node = OxmlElement(tag)
                rpr.append(node)
            node.set("typeface", font_face)
    except Exception:
        pass


def clamp_int(value: Any, minimum: int, maximum: int) -> int:
    try:
        number = int(float(value))
    except (TypeError, ValueError):
        return minimum
    return max(minimum, min(maximum, number))


def clamp_float(value: Any, minimum: float, maximum: float) -> float:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return minimum
    return max(minimum, min(maximum, number))


def write_pptx(path: Path, analysis: Dict[str, Any], source_image_path: Optional[Path] = None) -> None:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise RuntimeError("python-pptx is not installed. Run pip install -r api/requirements.txt") from exc

    prs = Presentation()
    slide_w, slide_h = slide_size_for_source(source_image_path)
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    editable_slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_editable_slide(editable_slide, analysis, source_image_path, slide_w, slide_h)

    include_source = os.getenv("PPT_INCLUDE_SOURCE_SLIDE", "0").strip().lower() in {"1", "true", "yes"}
    if source_image_path and include_source:
        visual_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_full_slide_image(visual_slide, source_image_path, slide_w, slide_h)

    temp_path = path.with_suffix(".tmp.pptx")
    prs.save(temp_path)
    os.replace(temp_path, path)


def slide_size_for_source(source_image_path: Optional[Path]) -> tuple[int, int]:
    if not source_image_path:
        return SLIDE_W, SLIDE_H
    try:
        from PIL import Image

        with Image.open(source_image_path) as image:
            image_w, image_h = image.size
        if image_w <= 0 or image_h <= 0:
            return SLIDE_W, SLIDE_H
        if image_w >= image_h:
            width = BASE_SLIDE_LONG_EDGE
            height = int(width * image_h / image_w)
        else:
            height = BASE_SLIDE_LONG_EDGE
            width = int(height * image_w / image_h)
        return max(3000000, width), max(3000000, height)
    except Exception:
        return SLIDE_W, SLIDE_H


def render_editable_slide(slide, analysis: Dict[str, Any], source_image_path: Optional[Path], slide_w: int, slide_h: int) -> None:
    from pptx.dml.color import RGBColor

    theme = analysis.get("theme") or {}
    background = clean_hex(theme.get("background")) or "10233F"
    accent = clean_hex(theme.get("accent")) or "1AA6D9"
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor.from_string(background)

    quality = normalize_quality(analysis.get("quality", "standard"))
    image_regions = analysis.get("image_regions") or []
    use_special_aeo_layout = should_use_aeo_layout(analysis)
    use_full_source_backdrop = (
        quality == "high_quality"
        and not use_special_aeo_layout
        and source_image_path
        and os.getenv("PPT_HQ_FULL_SOURCE_FALLBACK", "0").strip().lower() in {"1", "true", "yes"}
    )
    if use_full_source_backdrop:
        add_full_slide_image(slide, source_image_path, slide_w, slide_h)

    rendered_background_regions = False
    if use_full_source_backdrop:
        rendered_background_regions = True
    elif image_regions:
        render_image_regions(slide, analysis, source_image_path, layer="background", slide_w=slide_w, slide_h=slide_h)
        rendered_background_regions = True

    elements = build_dense_reconstruction(analysis) if use_special_aeo_layout else build_general_reconstruction(analysis)
    for element in elements:
        if element.get("type") == "imageRegions":
            layer = element.get("layer") or "background"
            if layer == "background" and rendered_background_regions:
                continue
            render_image_regions(slide, analysis, source_image_path, layer=layer, slide_w=slide_w, slide_h=slide_h)
            continue
        render_element(slide, element, accent, slide_w, slide_h)

    if image_regions and not use_full_source_backdrop:
        render_image_regions(slide, analysis, source_image_path, layer="foreground", slide_w=slide_w, slide_h=slide_h)


def should_use_aeo_layout(analysis: Dict[str, Any]) -> bool:
    # The old AEO process-sequence reconstruction was too specific and could leak
    # into unrelated AEO reports when AI analysis fell back. Keep all images on
    # the general reconstruction path unless a future dedicated detector is added.
    return False

def parse_step_heading(text: str):
    match = re.match(r"^\s*(\d{1,2})\s*[.:\-)]?\s+(.+?)\s*$", clean_text(text))
    if not match:
        return None
    label, heading = match.groups()
    try:
        number = int(label)
    except ValueError:
        return None
    if not 1 <= number <= 12:
        return None
    return str(number), heading


def enrich_high_quality_elements(analysis: Dict[str, Any], elements):
    enriched = list(elements)
    source_has_japanese = any(
        contains_japanese(clean_text(item.get("text")))
        for item in enriched
        if isinstance(item, dict)
    )
    if source_has_japanese:
        enriched = [
            item
            for item in enriched
            if not (
                isinstance(item, dict)
                and likely_non_source_english_analysis_text(clean_text(item.get("text")), source_has_japanese)
            )
        ]

    existing_text = compact_text(" ".join(clean_text(item.get("text")) for item in enriched if isinstance(item, dict)))
    theme = analysis.get("theme") or {}
    background = clean_hex(theme.get("background")) or "FFFFFF"
    font = readable_font_for_background(background)

    title = clean_text(analysis.get("title"))
    if title and not likely_non_source_english_analysis_text(title, source_has_japanese) and not text_already_present(title, existing_text):
        enriched.insert(0, {
            "type": "text",
            "text": title,
            "x": 0.055,
            "y": 0.045,
            "w": 0.62,
            "h": 0.075,
            "font": font,
            "font_face": clean_font_face("", title, 30),
            "font_size": 30,
            "bold": True,
        })
        existing_text += compact_text(title)

    subtitle = clean_text(analysis.get("subtitle"))
    if subtitle and not likely_non_source_english_analysis_text(subtitle, source_has_japanese) and not text_already_present(subtitle, existing_text):
        enriched.insert(1, {
            "type": "text",
            "text": subtitle,
            "x": 0.06,
            "y": 0.135,
            "w": 0.68,
            "h": 0.055,
            "font": font,
            "font_face": clean_font_face("", subtitle, 15),
            "font_size": 15,
            "bold": False,
        })
        existing_text += compact_text(subtitle)

    return enriched


def compact_text(text: str) -> str:
    return re.sub(r"\s+", "", text or "").lower()


def contains_japanese(text: str) -> bool:
    return bool(re.search(r"[\u3040-\u30ff\u3400-\u9fff]", text or ""))


def likely_non_source_english_analysis_text(text: str, source_has_japanese: bool) -> bool:
    if not source_has_japanese:
        return False
    cleaned = clean_text(text)
    if not cleaned or contains_japanese(cleaned):
        return False
    # Keep short visible acronyms/labels that commonly appear in Japanese designs.
    if re.fullmatch(r"[A-Z0-9+/#&.:\- ]{1,24}", cleaned):
        return False
    lowered = cleaned.lower()
    analysis_markers = (
        "hero section",
        "service strengths",
        "target pain",
        "construction type",
        "comparison",
        "cta",
        "footer",
        "main value proposition",
        "final conversion",
        "detailed list",
        "source image",
        "uploaded visual",
        "editable text",
        "generated as",
        "layout",
        "hierarchy",
    )
    if any(marker in lowered for marker in analysis_markers):
        return True
    english_words = re.findall(r"[A-Za-z]{3,}", cleaned)
    report_words = {
        "section",
        "summary",
        "proposition",
        "targeting",
        "detailed",
        "features",
        "service",
        "comparison",
        "conversion",
        "trust",
        "panels",
        "layout",
        "visual",
        "draft",
        "hierarchy",
        "editable",
    }
    return len(english_words) >= 4 and any(word.lower() in report_words for word in english_words)


def text_already_present(text: str, compact_existing: str) -> bool:
    compact = compact_text(text)
    if not compact:
        return True
    return bool(compact_existing) and (compact in compact_existing or compact_existing in compact)


def readable_font_for_background(background: str) -> str:
    try:
        red = int(background[0:2], 16)
        green = int(background[2:4], 16)
        blue = int(background[4:6], 16)
    except Exception:
        return "0B2341"
    luminance = 0.2126 * red + 0.7152 * green + 0.0722 * blue
    return "FFFFFF" if luminance < 105 else "0B2341"


def build_general_reconstruction(analysis: Dict[str, Any]):
    elements = []
    quality = normalize_quality(analysis.get("quality", "standard"))
    shape_transparency = 12 if quality == "high_quality" else 58
    for element in analysis.get("elements") or []:
        if not isinstance(element, dict):
            continue
        element_type = element.get("type")
        if element_type in {"text", "rect", "roundRect", "pill", "circle", "line"}:
            item = dict(element)
            if element_type not in {"text", "line"}:
                item.setdefault("transparency", shape_transparency)
            elements.append(item)
    if elements:
        if quality == "high_quality":
            elements = enrich_high_quality_elements(analysis, elements)
        return elements[:120 if quality == "high_quality" else 100]

    title = clean_text(analysis.get("title"))
    subtitle = clean_text(analysis.get("subtitle"))
    summary = clean_text(analysis.get("summary"))
    if title:
        elements.append({"type": "text", "text": title, "x": 0.06, "y": 0.05, "w": 0.62, "h": 0.08, "font": "0B2341", "font_size": 30, "bold": True})
    if subtitle:
        elements.append({"type": "text", "text": subtitle, "x": 0.06, "y": 0.15, "w": 0.62, "h": 0.06, "font": "1D3557", "font_size": 15, "bold": False})
    if summary:
        elements.append({"type": "text", "text": summary, "x": 0.06, "y": 0.84, "w": 0.88, "h": 0.06, "font": "1D3557", "font_size": 13, "bold": True})
    return elements


def render_image_regions(slide, analysis: Dict[str, Any], source_image_path: Optional[Path], layer: str, slide_w: int, slide_h: int) -> None:
    if not source_image_path or not source_image_path.exists():
        return

    regions = analysis.get("image_regions") or []
    for index, region in enumerate(regions[:MAX_IMAGE_REGIONS]):
        if not isinstance(region, dict):
            continue
        region_layer = clean_text(region.get("layer")).lower() or "foreground"
        if region_layer != layer:
            continue
        try:
            add_cropped_image_region(slide, source_image_path, region, index, slide_w, slide_h)
        except Exception as exc:
            print(
                f"Drop2PPT image region skipped index={index} error={redact_secret_text(str(exc))}",
                flush=True,
            )


def add_cropped_image_region(slide, source_image_path: Path, region: Dict[str, Any], index: int, slide_w: int, slide_h: int) -> None:
    from PIL import Image

    x_norm = clamp_float(region.get("x"), 0.0, 0.98)
    y_norm = clamp_float(region.get("y"), 0.0, 0.98)
    w_norm = clamp_float(region.get("w"), 0.01, 1.0)
    h_norm = clamp_float(region.get("h"), 0.01, 1.0)
    if x_norm + w_norm > 1.0:
        w_norm = max(0.01, 1.0 - x_norm)
    if y_norm + h_norm > 1.0:
        h_norm = max(0.01, 1.0 - y_norm)

    with Image.open(source_image_path) as image:
        image_w, image_h = image.size
        left = int(x_norm * image_w)
        top = int(y_norm * image_h)
        right = max(left + 2, int((x_norm + w_norm) * image_w))
        bottom = max(top + 2, int((y_norm + h_norm) * image_h))
        right = min(image_w, right)
        bottom = min(image_h, bottom)
        crop = image.crop((left, top, right, bottom))
        if crop.mode not in {"RGB", "RGBA"}:
            crop = crop.convert("RGBA")

        crop_dir = CROP_DIR / source_image_path.stem
        crop_dir.mkdir(parents=True, exist_ok=True)
        crop_path = crop_dir / f"region-{index}.png"
        crop.save(crop_path)

    slide_x = int(x_norm * slide_w)
    slide_y = int(y_norm * slide_h)
    region_w = int(w_norm * slide_w)
    region_h = int(h_norm * slide_h)
    slide.shapes.add_picture(str(crop_path), slide_x, slide_y, width=region_w, height=region_h)


def render_element(slide, element: Dict[str, Any], accent: str, slide_w: int = SLIDE_W, slide_h: int = SLIDE_H) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Pt

    x, y, w, h = element_box(element, slide_w, slide_h)
    element_type = element.get("type") or "rect"
    text = clean_text(element.get("text"))
    fill = clean_hex(element.get("fill")) or "123B59"
    line = clean_hex(element.get("line")) or accent
    font = clean_hex(element.get("font")) or "FFFFFF"
    font_size = clamp_int(element.get("font_size"), 7, 44)
    font_face = clean_font_face(element.get("font_face") or element.get("font_family") or element.get("typeface"), text, font_size)
    bold = bool(element.get("bold"))
    transparency = clamp_int(element.get("transparency"), 0, 100)

    if element_type == "line":
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y + h)
        connector.line.color.rgb = RGBColor.from_string(line)
        connector.line.width = 18000
        return

    if element_type == "text":
        add_textbox(slide, x, y, w, h, text, font, font_size, bold=bold, font_face=font_face)
        return

    shape_type = {
        "circle": MSO_SHAPE.OVAL,
        "pill": MSO_SHAPE.ROUNDED_RECTANGLE,
        "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "rect": MSO_SHAPE.RECTANGLE,
    }.get(element_type, MSO_SHAPE.ROUNDED_RECTANGLE)

    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    if transparency:
        shape.fill.transparency = transparency
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = 15000

    if text:
        frame = shape.text_frame
        frame.clear()
        frame.margin_left = 70000
        frame.margin_right = 70000
        frame.margin_top = 50000
        frame.margin_bottom = 50000
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        for index, line_text in enumerate(text.splitlines()[:5]):
            p = paragraph if index == 0 else frame.add_paragraph()
            run = p.add_run()
            run.text = line_text
            run.font.bold = bold or index == 0
            run.font.size = Pt(font_size if index == 0 else max(7, font_size - 2))
            run.font.color.rgb = RGBColor.from_string(font)
            apply_run_font(run, font_face)


def element_box(element: Dict[str, Any], slide_w: int = SLIDE_W, slide_h: int = SLIDE_H):
    is_line = (element.get("type") or "").lower() == "line"
    x = int(clamp_float(element.get("x"), 0.0, 0.98) * slide_w)
    y = int(clamp_float(element.get("y"), 0.0, 0.98) * slide_h)
    w = int(clamp_float(element.get("w"), 0.0 if is_line else 0.01, 1.0) * slide_w)
    h = int(clamp_float(element.get("h"), 0.0 if is_line else 0.01, 1.0) * slide_h)
    return x, y, w, h


def add_full_slide_image(slide, image_path: Path, slide_w: int = SLIDE_W, slide_h: int = SLIDE_H) -> None:
    try:
        from PIL import Image

        with Image.open(image_path) as image:
            image_w, image_h = image.size
        image_ratio = image_w / image_h
        slide_ratio = slide_w / slide_h
        if image_ratio >= slide_ratio:
            width = slide_w
            height = int(slide_w / image_ratio)
        else:
            height = slide_h
            width = int(slide_h * image_ratio)
        left = int((slide_w - width) / 2)
        top = int((slide_h - height) / 2)
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    except Exception:
        slide.shapes.add_picture(str(image_path), 0, 0, width=slide_w, height=slide_h)


def add_textbox(slide, x: int, y: int, w: int, h: int, text: str, font: str, font_size: int, bold: bool = False, font_face: str = "") -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor.from_string(font)
    apply_run_font(run, clean_font_face(font_face, text, font_size))


def add_block(
    slide,
    shape_type,
    x: int,
    y: int,
    w: int,
    h: int,
    title: str,
    body: str,
    fill: str = "123B59",
    line: str = "1AA6D9",
    font: str = "FFFFFF",
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = 18000

    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 120000
    frame.margin_right = 120000
    frame.margin_top = 90000
    frame.margin_bottom = 90000
    frame.word_wrap = True

    title_p = frame.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = title
    title_run.font.bold = True
    title_run.font.size = Pt(14)
    title_run.font.color.rgb = RGBColor.from_string(font)
    apply_run_font(title_run, clean_font_face("", title, 14))

    if body:
        body_p = frame.add_paragraph()
        body_run = body_p.add_run()
        body_run.text = body
        body_run.font.size = Pt(11)
        body_run.font.color.rgb = RGBColor.from_string(font)
        apply_run_font(body_run, clean_font_face("", body, 11))


def build_slide_xml(analysis: Dict[str, Any]) -> str:
    shapes = [
        shape_xml(2, 0, 0, SLIDE_W, SLIDE_H, "", "", fill="10233F", line="10233F", font="FFFFFF", font_size=1200),
        textbox_xml(3, 420000, 260000, 8300000, 700000, analysis["title"], "FFFFFF", 3400, bold=True),
        textbox_xml(4, 470000, 940000, 8000000, 360000, analysis["subtitle"], "D9F3FF", 1350, bold=False),
        textbox_xml(5, 500000, 4480000, 8100000, 360000, analysis["summary"], "EAF7FF", 1100, bold=False),
    ]
    shape_id = 10
    for section in analysis["sections"]:
        x = int(section["x"] * SLIDE_W)
        y = int(section["y"] * SLIDE_H)
        w = int(section["w"] * SLIDE_W)
        h = int(section["h"] * SLIDE_H)
        shapes.append(shape_xml(shape_id, x, y, w, h, section["title"], section["body"]))
        shape_id += 1

    step_w = int(SLIDE_W * 0.20)
    gap = int(SLIDE_W * 0.02)
    start_x = int(SLIDE_W * 0.08)
    y = int(SLIDE_H * 0.76)
    for index, step in enumerate(analysis["steps"][:4]):
        x = start_x + index * (step_w + gap)
        title = f"{step['label']}  {step['title']}"
        shapes.append(shape_xml(shape_id, x, y, step_w, int(SLIDE_H * 0.14), title, step["body"], fill="E9F7F7", line="1AA6D9", font="12324A"))
        shape_id += 1

    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/>
          <a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      {''.join(shapes)}
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>"""


def shape_xml(shape_id: int, x: int, y: int, w: int, h: int, title: str, body: str, fill: str = "123B59", line: str = "1AA6D9", font: str = "FFFFFF", font_size: int = 1050) -> str:
    text = ""
    if title or body:
        text = f"""
      <p:txBody>
        <a:bodyPr wrap="square" lIns="120000" tIns="90000" rIns="120000" bIns="90000"/>
        <a:lstStyle/>
        <a:p><a:r><a:rPr lang="en-US" sz="{font_size + 250}" b="1"><a:solidFill><a:srgbClr val="{font}"/></a:solidFill></a:rPr><a:t>{escape(title)}</a:t></a:r></a:p>
        <a:p><a:r><a:rPr lang="en-US" sz="{font_size}"><a:solidFill><a:srgbClr val="{font}"/></a:solidFill></a:rPr><a:t>{escape(body)}</a:t></a:r></a:p>
      </p:txBody>"""
    return f"""
      <p:sp>
        <p:nvSpPr><p:cNvPr id="{shape_id}" name="Editable block {shape_id}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
        <p:spPr>
          <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>
          <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
          <a:solidFill><a:srgbClr val="{fill}"><a:alpha val="90000"/></a:srgbClr></a:solidFill>
          <a:ln w="18000"><a:solidFill><a:srgbClr val="{line}"/></a:solidFill></a:ln>
        </p:spPr>
        {text}
      </p:sp>"""


def textbox_xml(shape_id: int, x: int, y: int, w: int, h: int, text: str, font: str, font_size: int, bold: bool = False) -> str:
    return f"""
      <p:sp>
        <p:nvSpPr><p:cNvPr id="{shape_id}" name="Editable text {shape_id}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
        <p:spPr>
          <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>
          <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
          <a:noFill/>
        </p:spPr>
        <p:txBody>
          <a:bodyPr wrap="square"/>
          <a:lstStyle/>
          <a:p><a:r><a:rPr lang="en-US" sz="{font_size}" b="{1 if bold else 0}"><a:solidFill><a:srgbClr val="{font}"/></a:solidFill></a:rPr><a:t>{escape(text)}</a:t></a:r></a:p>
        </p:txBody>
      </p:sp>"""


def content_types_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
  <Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
  <Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
  <Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
  <Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
  <Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
</Types>"""


def package_rels_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>
</Relationships>"""


def presentation_xml() -> str:
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>
  <p:sldIdLst><p:sldId id="256" r:id="rId2"/></p:sldIdLst>
  <p:sldSz cx="{SLIDE_W}" cy="{SLIDE_H}" type="screen16x9"/>
  <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>"""


def presentation_rels_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="theme/theme1.xml"/>
</Relationships>"""


def empty_rels_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>"""


def slide_master_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld>
  <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
  <p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst>
  <p:txStyles><p:titleStyle/><p:bodyStyle/><p:otherStyle/></p:txStyles>
</p:sldMaster>"""


def slide_master_rels_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>"""


def slide_layout_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank" preserve="1">
  <p:cSld name="Blank"><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sldLayout>"""


def theme_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="Drop2PPT">
  <a:themeElements>
    <a:clrScheme name="Drop2PPT"><a:dk1><a:srgbClr val="14212F"/></a:dk1><a:lt1><a:srgbClr val="FFFFFF"/></a:lt1><a:dk2><a:srgbClr val="12324A"/></a:dk2><a:lt2><a:srgbClr val="F2F6F8"/></a:lt2><a:accent1><a:srgbClr val="0F8D80"/></a:accent1><a:accent2><a:srgbClr val="1AA6D9"/></a:accent2><a:accent3><a:srgbClr val="D09B24"/></a:accent3><a:accent4><a:srgbClr val="D35C43"/></a:accent4><a:accent5><a:srgbClr val="516173"/></a:accent5><a:accent6><a:srgbClr val="10233F"/></a:accent6><a:hlink><a:srgbClr val="1AA6D9"/></a:hlink><a:folHlink><a:srgbClr val="0F8D80"/></a:folHlink></a:clrScheme>
    <a:fontScheme name="Office"><a:majorFont><a:latin typeface="Aptos Display"/></a:majorFont><a:minorFont><a:latin typeface="Aptos"/></a:minorFont></a:fontScheme>
    <a:fmtScheme name="Office"><a:fillStyleLst/><a:lnStyleLst/><a:effectStyleLst/><a:bgFillStyleLst/></a:fmtScheme>
  </a:themeElements>
</a:theme>"""


def core_props_xml() -> str:
    now = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <dc:title>Drop2PPT editable slide</dc:title>
  <dc:creator>WorldScene</dc:creator>
  <cp:lastModifiedBy>WorldScene</cp:lastModifiedBy>
  <dcterms:created xsi:type="dcterms:W3CDTF">{now}</dcterms:created>
  <dcterms:modified xsi:type="dcterms:W3CDTF">{now}</dcterms:modified>
</cp:coreProperties>"""


def app_props_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">
  <Application>WorldScene Drop2PPT</Application>
  <PresentationFormat>On-screen Show (16:9)</PresentationFormat>
  <Slides>1</Slides>
</Properties>"""
